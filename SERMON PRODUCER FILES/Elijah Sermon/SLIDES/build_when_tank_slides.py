from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "When_the_Tank_Is_Not_Low_Its_Empty_sermon_slides.pptx"

SLIDES = [
    {
        "kind": "cover",
        "kicker": "1 Kings 19:1-18",
        "title": "When the Tank Is Not Low.\nIt's Empty.",
        "subtitle": "What God does when you have nothing left",
    },
    {
        "kind": "quote",
        "kicker": "Central Proposition",
        "lines": ["At empty, God's first response is not correction or commission.", "It is care."],
    },
    {"kind": "quote", "lines": ["Chapter 19 is not about fire.", "It is about bread."]},
    {
        "kind": "quote",
        "kicker": "From Mount Carmel to the wilderness",
        "lines": ["The distance between your greatest victory and your deepest collapse can be", "one bad message."],
    },
    {"kind": "point", "kicker": "Point One", "title": "God meets us at empty,\nnot after empty."},
    {
        "kind": "two_col",
        "left_title": "What God Does Not Do",
        "left": ["Rebuke", "Replay", "Rush", "Assign"],
        "right_title": "What God Does",
        "right": ["Bread", "Water", "Sleep", "Touch"],
    },
    {"kind": "quote", "kicker": "1 Kings 19:7", "lines": ['"The journey is too great for thee."']},
    {
        "kind": "quote",
        "lines": ["Sabbath is the creation architecture of rest.", "Under the juniper tree, God deploys Sabbath as mercy."],
    },
    {"kind": "point", "kicker": "Point Two", "title": "The God who speaks\nin the silence."},
    {"kind": "contrast", "left": "Wind\nEarthquake\nFire", "right": "Not there."},
    {
        "kind": "quote",
        "lines": ["God made noise to show Elijah He could.", "Then He chose silence to show Elijah who He was."],
    },
    {"kind": "quote", "lines": ["What noise is making it impossible to hear the thin silence?"]},
    {"kind": "point", "kicker": "Point Three", "title": "God commissions\nbroken people."},
    {"kind": "quote", "lines": ["Grace does not require full restoration before it trusts you again."]},
    {
        "kind": "quote",
        "kicker": "1 Kings 19:18",
        "lines": ["There are still seven thousand.", "God's work is never hanging by one person's endurance."],
    },
    {
        "kind": "closing",
        "title": "The bread is already baked.",
        "subtitle": '"Arise and eat; because the journey is too great for thee."',
    },
]


COLORS = {
    "night": RGBColor(23, 27, 31),
    "charcoal": RGBColor(34, 39, 44),
    "cream": RGBColor(245, 239, 226),
    "bread": RGBColor(215, 168, 79),
    "ember": RGBColor(185, 90, 61),
    "sage": RGBColor(143, 163, 143),
    "mist": RGBColor(203, 208, 195),
}


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["night"]

    wash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    wash.fill.solid()
    wash.fill.fore_color.rgb = COLORS["charcoal"]
    wash.fill.transparency = 28
    wash.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.95), Inches(5.6), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["bread"]
    accent.line.fill.background()

    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.4), Inches(-1.2), Inches(4.2), Inches(4.2))
    halo.fill.solid()
    halo.fill.fore_color.rgb = COLORS["bread"]
    halo.fill.transparency = 82
    halo.line.fill.background()


def textbox(slide, text, x, y, w, h, size, color="cream", bold=False, serif=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Georgia" if serif else "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = COLORS[color]
    return shape


def add_kicker(slide, text):
    textbox(slide, text.upper(), 0.76, 0.58, 7.2, 0.38, 15, "mist", bold=True)


def cover(slide, data):
    add_bg(slide)
    add_kicker(slide, data["kicker"])
    textbox(slide, data["title"], 0.72, 1.35, 10.2, 2.35, 46, "cream", bold=True, serif=True)
    textbox(slide, data["subtitle"], 0.78, 4.05, 7.6, 0.52, 22, "mist")
    textbox(slide, "Braselton Seventh-day Adventist Church", 0.78, 6.42, 6.2, 0.3, 13, "sage")


def quote(slide, data):
    add_bg(slide)
    if data.get("kicker"):
        add_kicker(slide, data["kicker"])
    y = 1.35 if data.get("kicker") else 1.58
    lines = data["lines"]
    first_size = 39 if len(lines[0]) > 70 else 43
    textbox(slide, lines[0], 0.78, y, 10.7, 1.75, first_size, "cream", serif=True)
    if len(lines) > 1:
        second_size = 40 if len(lines[1]) > 68 else 45
        textbox(slide, lines[1], 0.78, y + 2.0, 10.6, 1.25, second_size, "bread", bold=True, serif=True)


def point(slide, data):
    add_bg(slide)
    add_kicker(slide, data["kicker"])
    textbox(slide, data["title"], 0.78, 1.72, 10.7, 2.7, 52, "cream", bold=True, serif=True)


def two_col(slide, data):
    add_bg(slide)
    textbox(slide, data["left_title"].upper(), 0.78, 0.72, 5.1, 0.36, 16, "bread", bold=True)
    textbox(slide, data["right_title"].upper(), 7.02, 0.72, 5.1, 0.36, 16, "bread", bold=True)
    for idx, word in enumerate(data["left"]):
        y = 1.55 + idx * 1.0
        bullet(slide, 0.92, y + 0.13)
        textbox(slide, word, 1.18, y, 4.3, 0.52, 31, "cream", serif=True)
    for idx, word in enumerate(data["right"]):
        y = 1.55 + idx * 1.0
        bullet(slide, 7.16, y + 0.13)
        textbox(slide, word, 7.42, y, 4.3, 0.52, 31, "cream", serif=True)


def bullet(slide, x, y):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLORS["bread"]
    dot.line.fill.background()


def contrast(slide, data):
    add_bg(slide)
    textbox(slide, data["left"], 0.8, 1.38, 5.2, 3.8, 50, "mist", bold=True, serif=True)
    textbox(slide, data["right"], 6.8, 2.35, 5.0, 1.1, 57, "bread", bold=True, serif=True)


def closing(slide, data):
    add_bg(slide)
    textbox(slide, data["title"], 0.78, 2.08, 9.7, 1.18, 49, "cream", bold=True, serif=True)
    textbox(slide, data["subtitle"], 0.82, 3.62, 9.8, 0.82, 24, "mist", serif=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    renderers = {
        "cover": cover,
        "quote": quote,
        "point": point,
        "two_col": two_col,
        "contrast": contrast,
        "closing": closing,
    }
    for data in SLIDES:
        slide = prs.slides.add_slide(blank)
        renderers[data["kind"]](slide, data)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
