from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
PREVIEW_DIR = ROOT / "preview_png"
OUT = ROOT / "preview_contact_sheet.png"
PPTX = ROOT / "When_the_Tank_Is_Not_Low_Its_Empty_sermon_slides.pptx"


def build_contact_sheet():
    files = sorted(PREVIEW_DIR.glob("slide.*.png"))
    thumbs = []
    for idx, file in enumerate(files, start=1):
        img = Image.open(file).convert("RGB")
        img.thumbnail((320, 180))
        canvas = Image.new("RGB", (340, 220), "white")
        canvas.paste(img, (10, 10))
        draw = ImageDraw.Draw(canvas)
        draw.text((12, 194), f"Slide {idx:02d}", fill=(20, 20, 20))
        thumbs.append(canvas)

    cols = 4
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 340, rows * 220), (245, 245, 245))
    for i, thumb in enumerate(thumbs):
        x = (i % cols) * 340
        y = (i // cols) * 220
        sheet.paste(thumb, (x, y))
    sheet.save(OUT)
    return len(files)


def inspect_pptx():
    prs = Presentation(PPTX)
    text_shapes = 0
    empty_text_shapes = 0
    sample_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_shapes += 1
                text = shape.text.strip()
                if text:
                    sample_text.append(text)
                else:
                    empty_text_shapes += 1
    return len(prs.slides), text_shapes, empty_text_shapes, sample_text[:8]


if __name__ == "__main__":
    previews = build_contact_sheet()
    slide_count, text_shapes, empty_text_shapes, sample_text = inspect_pptx()
    print(f"preview_count={previews}")
    print(f"pptx_slide_count={slide_count}")
    print(f"pptx_text_shapes={text_shapes}")
    print(f"pptx_empty_text_shapes={empty_text_shapes}")
    print("pptx_sample_text=" + " | ".join(sample_text))
    print(f"contact_sheet={OUT}")
